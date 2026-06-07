import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Initialize Presentation and set to 16:9 widescreen
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color Palette
DARK_GREEN = RGBColor(14, 77, 68)      # #0E4D44 (Primary)
MINT_TEAL = RGBColor(26, 163, 132)     # #1AA384 (Secondary)
LIGHT_MINT = RGBColor(143, 232, 210)   # #8FE8D2 (Accent)
ALERT_RED = RGBColor(192, 57, 43)      # #C0392B (Alert)
TEXT_DARK = RGBColor(33, 33, 33)       # #212121 (Primary Body Text)
BG_LIGHT = RGBColor(246, 248, 246)     # #F6F8F6 (Light Background)
WHITE = RGBColor(255, 255, 255)

# Font configurations
FONT_TITLE = "Helvetica"
FONT_BODY = "Arial"

def set_slide_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def create_title_textbox(slide, text, color=DARK_GREEN, size=Pt(36)):
    txBox = slide.shapes.add_textbox(Inches(0.75), Inches(0.4), Inches(11.83), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_top = tf.margin_bottom = tf.margin_left = tf.margin_right = 0
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = FONT_TITLE
    p.font.size = size
    p.font.bold = True
    p.font.color.rgb = color
    return txBox

def add_footer(slide, current_slide, total_slides=7):
    # Left Footer: Project title
    txBox = slide.shapes.add_textbox(Inches(0.75), Inches(7.0), Inches(8.0), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Open Project 2026 | Society of Business, IIT Roorkee"
    p.font.name = FONT_BODY
    p.font.size = Pt(9)
    p.font.color.rgb = MINT_TEAL
    
    # Right Footer: Page number
    txBox_num = slide.shapes.add_textbox(Inches(11.53), Inches(7.0), Inches(1.05), Inches(0.3))
    tf_num = txBox_num.text_frame
    p_num = tf_num.paragraphs[0]
    p_num.alignment = PP_ALIGN.RIGHT
    p_num.text = f"Slide {current_slide} of {total_slides}"
    p_num.font.name = FONT_BODY
    p_num.font.size = Pt(9)
    p_num.font.color.rgb = MINT_TEAL

# Slide 0: Cover Slide (Dark Theme)
blank_slide_layout = prs.slide_layouts[6]
slide0 = prs.slides.add_slide(blank_slide_layout)
set_slide_background(slide0, DARK_GREEN)

# Large Title Box
txBox = slide0.shapes.add_textbox(Inches(0.75), Inches(2.2), Inches(11.83), Inches(2.5))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Agentic AI-Based Dynamic Tariff\nOptimization for EV Charging Networks"
p.font.name = FONT_TITLE
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = WHITE

# Subtitle
p2 = tf.add_paragraph()
p2.text = "\nLeveraging Large-Scale Charging Session Data (Caltech/JPL & UrbanEV)"
p2.font.name = FONT_BODY
p2.font.size = Pt(18)
p2.font.color.rgb = LIGHT_MINT

# Author / Institution
txBox_auth = slide0.shapes.add_textbox(Inches(0.75), Inches(5.2), Inches(11.83), Inches(1.2))
tf_auth = txBox_auth.text_frame
p3 = tf_auth.paragraphs[0]
p3.text = "Open Project 2026 — Society of Business, IIT Roorkee"
p3.font.name = FONT_BODY
p3.font.size = Pt(14)
p3.font.bold = True
p3.font.color.rgb = WHITE

p4 = tf_auth.add_paragraph()
p4.text = "Case Study Submission | Target Deadline: June 7th, 2026"
p4.font.name = FONT_BODY
p4.font.size = Pt(12)
p4.font.color.rgb = LIGHT_MINT


# Slide 1: Data Landscape & Preprocessing (Light Theme)
slide1 = prs.slides.add_slide(blank_slide_layout)
set_slide_background(slide1, BG_LIGHT)
create_title_textbox(slide1, "Data Landscape & Preprocessing Decisions")
add_footer(slide1, 1)

# Two Column Text Blocks
# Left Column: Data Sources
tx_left = slide1.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(5.6), Inches(5.0))
tf_left = tx_left.text_frame
tf_left.word_wrap = True
p = tf_left.paragraphs[0]
p.text = "Data Landscape Overview"
p.font.name = FONT_TITLE
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = DARK_GREEN

bullets_left = [
    "ACN-Data (Caltech/JPL): 14,199 charging sessions representing US workplace patterns. Contains connection/disconnection timestamps, user energy requests (kWh), and done-charging timings.",
    "UrbanEV (Shenzhen): Time-series occupancy matrix for 24,798 charging piles in 5-minute intervals. Captured across residential and CBD grids.",
    "Geographic Mismatch Strategy: Datasets are geographically disjoint. ACN provides session behavior distributions (charging-to-session ratios); UrbanEV provides occupancy signals for predictive modeling."
]
for bullet in bullets_left:
    p = tf_left.add_paragraph()
    p.text = "• " + bullet
    p.font.name = FONT_BODY
    p.font.size = Pt(13)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(10)

# Right Column: Preprocessing Decisions
tx_right = slide1.shapes.add_textbox(Inches(6.9), Inches(1.5), Inches(5.6), Inches(5.0))
tf_right = tx_right.text_frame
tf_right.word_wrap = True
p = tf_right.paragraphs[0]
p.text = "Feature Engineering & Imputations"
p.font.name = FONT_TITLE
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = DARK_GREEN

bullets_right = [
    "Queue Length Proxy: Waiting times are modeled when occupancy rate exceeds 85%: Queue = (Occupancy - 0.85) * Capacity * 1.5.",
    "Grid Cost Tiers: 3-tier Time-of-Day (ToD) tariff structure based on Indian DISCOMs (₹4 off-peak, ₹7 shoulder, ₹10 peak procurement).",
    "Outlier Capping: Session durations clipped at the 99th percentile to exclude vehicles parked indefinitely without charging.",
    "Missing Values: Dropped ACN sessions with zero or negative durations; missing ACN doneChargingTime imputed as disconnectTime."
]
for bullet in bullets_right:
    p = tf_right.add_paragraph()
    p.text = "• " + bullet
    p.font.name = FONT_BODY
    p.font.size = Pt(13)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(10)


# Slide 2: EDA & Demand Insights (Light Theme)
slide2 = prs.slides.add_slide(blank_slide_layout)
set_slide_background(slide2, BG_LIGHT)
create_title_textbox(slide2, "Exploratory Data Analysis & Demand Profiles")
add_footer(slide2, 2)

# Left Column Text
tx = slide2.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(5.2), Inches(5.0))
tf = tx.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Key Behavioral Insights"
p.font.name = FONT_TITLE
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = DARK_GREEN

bullets = [
    "CBD vs. Residential Spikes: CBD zones show sharp, high occupancy spikes during working hours (9:00 - 17:00), frequently exceeding the 80% congestion threshold. Residential zones show flatter, evening-biased demand.",
    "Weekday vs. Weekend Demand: Commuter traffic creates clear weekday peak windows, while weekends show lower average demand without sharp spikes.",
    "User Type Signatures: Fleet/Repeat users represent a small user base but charge highly predictably and draw larger energy portions (Avg: 14.05 kWh) compared to casual public users (Avg: 8.45 kWh)."
]
for bullet in bullets:
    p = tf.add_paragraph()
    p.text = "• " + bullet
    p.font.name = FONT_BODY
    p.font.size = Pt(13)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(10)

# Right Column: Embed Plot
plot_path = "/Users/anubhav/Documents/EV_Tariff_Optimization/plots/1_cbd_vs_res_demand.png"
if os.path.exists(plot_path):
    slide2.shapes.add_picture(plot_path, Inches(6.3), Inches(1.6), Inches(6.5), Inches(2.32))

plot_path_fleet = "/Users/anubhav/Documents/EV_Tariff_Optimization/plots/6_fleet_vs_public.png"
if os.path.exists(plot_path_fleet):
    slide2.shapes.add_picture(plot_path_fleet, Inches(6.3), Inches(4.1), Inches(6.5), Inches(2.5))


# Slide 3: Demand Prediction Agent (Light Theme)
slide3 = prs.slides.add_slide(blank_slide_layout)
set_slide_background(slide3, BG_LIGHT)
create_title_textbox(slide3, "Agent 1: ML-Based Demand Prediction Agent")
add_footer(slide3, 3)

# Left Column Text
tx = slide3.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(5.2), Inches(5.0))
tf = tx.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Model Architecture & Performance"
p.font.name = FONT_TITLE
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = DARK_GREEN

bullets = [
    "Predictive Model: Gradient Boosting Regressor (GBR) chosen for its ability to map non-linear plateaus and sharp transitions in occupancy.",
    "Hyperparameters: n_estimators=100, max_depth=5. Capped depth prevents overfitting on highly autocorrelated lag features.",
    "Features: Short-term lags (1, 3, 6, 12 steps), rolling averages, temporal flags (hour, day of week), spatial type (is_cbd), and grid costs.",
    "Model Accuracy on Test Set:\n  - RMSE: 0.0129 (1.29% error in occupancy rate)\n  - MAE: 0.0048\n  - R² Score: 0.9947 (99.47% variance explained)",
    "Outputs: Occupancy forecasts, sigmoid-scaled congestion probabilities, and Expected Charging Load (Total: 83.767 GWh)."
]
for bullet in bullets:
    p = tf.add_paragraph()
    p.text = "• " + bullet if not bullet.startswith("  -") else bullet
    p.font.name = FONT_BODY
    p.font.size = Pt(13)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(8)

# Right Column: Embed Plot
plot_path = "/Users/anubhav/Documents/EV_Tariff_Optimization/plots/8_demand_prediction_diagnostics.png"
if os.path.exists(plot_path):
    slide3.shapes.add_picture(plot_path, Inches(6.3), Inches(1.8), Inches(6.5), Inches(4.3))


# Slide 4: Tariff Pricing Agent (Light Theme)
slide4 = prs.slides.add_slide(blank_slide_layout)
set_slide_background(slide4, BG_LIGHT)
create_title_textbox(slide4, "Agent 2: Tariff Pricing Agent & Elasticity")
add_footer(slide4, 4)

# Left Column Text
tx = slide4.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(5.2), Inches(5.0))
tf = tx.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Dynamic Pricing Rules & Elasticity"
p.font.name = FONT_TITLE
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = DARK_GREEN

bullets = [
    "Tariff Rules (₹15/kWh flat baseline):\n  - Surge Pricing (>80% occupancy): up to +50% (₹22.50)\n  - Discount Pricing (<30% occupancy): down to -30% (₹10.50)",
    "Asymmetric Customer Elasticity Model:\n  - Peak Surcharges (Elasticity = -0.10): Users are inelastic. High captive demand, low charging flexibility.\n  - Off-Peak Discounts (Elasticity = -0.85): Users are highly elastic. Incentivizes shift in discretionary charging times.",
    "Peak Shaving & Valley Filling: Shaves peak congestion and redirects 45% of peak-shaved load into off-peak valley slots."
]
for bullet in bullets:
    p = tf.add_paragraph()
    p.text = "• " + bullet if not bullet.startswith("  -") else bullet
    p.font.name = FONT_BODY
    p.font.size = Pt(13)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(10)

# Right Column: Embed Plot
plot_path = "/Users/anubhav/Documents/EV_Tariff_Optimization/plots/9_peak_shaving_impact.png"
if os.path.exists(plot_path):
    slide4.shapes.add_picture(plot_path, Inches(6.3), Inches(1.5), Inches(6.5), Inches(2.95))

plot_path_grid = "/Users/anubhav/Documents/EV_Tariff_Optimization/plots/10_dynamic_tariff_vs_grid.png"
if os.path.exists(plot_path_grid):
    slide4.shapes.add_picture(plot_path_grid, Inches(6.3), Inches(4.5), Inches(6.5), Inches(2.2))


# Slide 5: Monitoring & Learning Agent (Light Theme)
slide5 = prs.slides.add_slide(blank_slide_layout)
set_slide_background(slide5, BG_LIGHT)
create_title_textbox(slide5, "Agent 3: Monitoring & Threshold Optimization")
add_footer(slide5, 5)

# Left Column Text
tx = slide5.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(5.2), Inches(5.0))
tf = tx.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Feedback Loop & Threshold Convergence"
p.font.name = FONT_TITLE
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = DARK_GREEN

bullets = [
    "Feedback Architecture: Evaluates policy rewards over successive daily episodes. Optimizes thresholds using local coordinate ascent (step size = 0.015).",
    "Multi-Objective Reward Function:\n  Reward = Revenue Gain % + 0.10 * Queue Reduction %",
    "Optimization Convergence (Episode 40):\n  - Peak Threshold: Converges from 0.800 to 0.300\n  - Off-Peak Threshold: Converges from 0.300 to 0.050",
    "Commercial Logic: Because peak users are highly price inelastic, lowering the peak threshold to 0.300 maximizes surcharged hours and revenue gain (+6.75%), while restricting discounts to preserve margin."
]
for bullet in bullets:
    p = tf.add_paragraph()
    p.text = "• " + bullet if not bullet.startswith("  -") and not bullet.startswith("  Reward") else bullet
    p.font.name = FONT_BODY
    p.font.size = Pt(13)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(8)

# Right Column: Embed Plot
plot_path = "/Users/anubhav/Documents/EV_Tariff_Optimization/plots/11_threshold_optimization_curve.png"
if os.path.exists(plot_path):
    slide5.shapes.add_picture(plot_path, Inches(6.3), Inches(1.5), Inches(6.5), Inches(2.75))

plot_path_conv = "/Users/anubhav/Documents/EV_Tariff_Optimization/plots/13_threshold_convergence.png"
if os.path.exists(plot_path_conv):
    slide5.shapes.add_picture(plot_path_conv, Inches(6.3), Inches(4.35), Inches(6.5), Inches(2.4))


# Slide 6: Implications & Results (Light Theme)
slide6 = prs.slides.add_slide(blank_slide_layout)
set_slide_background(slide6, BG_LIGHT)
create_title_textbox(slide6, "Project Outcomes & Implications")
add_footer(slide6, 6)

# Left Column Text
tx = slide6.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(5.2), Inches(5.0))
tf = tx.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Key Results & Impact Summary"
p.font.name = FONT_TITLE
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = DARK_GREEN

bullets = [
    "Revenue Gain: +6.75% overall increase in operator revenue. Driven by an increase in average pricing yield to 16.126 ₹/kWh (from a flat baseline of ₹15.00).",
    "Congestion Relief: -53.77% reduction in peak-hour wait lines and queue lengths. Driven by peak shaving.",
    "Off-Peak Uplift: +8.23% increase in sessions and occupancy during off-peak valleys, filling idle capacity.",
    "Customer Response Rate: +2.95% overall session volume shift.",
    "Grid & Policy Benefits: Flattening the charging profile protects grid transformer components, reduces peak network load, and supports lower-cost public charging."
]
for bullet in bullets:
    p = tf.add_paragraph()
    p.text = "• " + bullet
    p.font.name = FONT_BODY
    p.font.size = Pt(13)
    p.font.color.rgb = TEXT_DARK
    p.space_before = Pt(8)

# Right Column: Embed Plot
plot_path = "/Users/anubhav/Documents/EV_Tariff_Optimization/plots/14_kpi_comparison.png"
if os.path.exists(plot_path):
    slide6.shapes.add_picture(plot_path, Inches(6.3), Inches(1.8), Inches(6.5), Inches(4.3))


# Slide 7: Appendix & Future Scope (Dark Theme)
slide7 = prs.slides.add_slide(blank_slide_layout)
set_slide_background(slide7, DARK_GREEN)

# Title in White
create_title_textbox(slide7, "Appendix, Limitations & Future Scope", color=WHITE, size=Pt(32))
add_footer(slide7, 7)

# Text Box
tx = slide7.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(11.83), Inches(5.0))
tf = tx.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "Project Limitations & Scope Constraints"
p.font.name = FONT_TITLE
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = LIGHT_MINT

bullets_app = [
    "Geographic Mismatch: ACN-Data (US sessions) and UrbanEV (Shenzhen zones) requires assuming user session shapes generalize across geographies. A unified deployment requires local station arrival logs.",
    "Offline Coordinate Search vs. Online Reinforcement Learning: Coordinate ascent is optimized over a static historical dataset. A production system must implement online RL (such as Contextual Bandits) on stream observations.",
    "Elasticity Calibration: Customer elasticities (-0.10 surcharges, -0.85 discounts) are literature-based assumptions. Real deployments require live A/B tests to estimate price elasticity curves.",
    "Future Directions: Co-optimize pricing tariffs with real-time solar/wind availability to maximize green charging. Extend agents to support Vehicle-to-Grid (V2G) discharging incentives."
]
for bullet in bullets_app:
    p = tf.add_paragraph()
    p.text = "• " + bullet
    p.font.name = FONT_BODY
    p.font.size = Pt(13)
    p.font.color.rgb = WHITE
    p.space_before = Pt(12)

# Save presentation
prs.save("/Users/anubhav/Documents/EV_Tariff_Optimization/EV_Tariff_Optimization_Presentation.pptx")
print("Presentation deck generated successfully as 'EV_Tariff_Optimization_Presentation.pptx'.")
